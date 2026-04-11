from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x39, 0x64)   # title bg
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # section header bg
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)   # accent / table header
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)   # highlight numbers

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper: add text box ─────────────────────────────────────────────────────
def tb(slide, text, l, t, w, h, size=18, bold=False, color=DARK_GREY,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def tb_lines(slide, lines, l, t, w, h, size=16, color=DARK_GREY,
             bold_first=False, line_space=None):
    """Add a textbox with multiple paragraphs."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if line_space:
            p.space_before = Pt(line_space)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb

def filled_rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ── Slide builder helpers ────────────────────────────────────────────────────
def title_slide(title, subtitle=""):
    sld = prs.slides.add_slide(BLANK)
    filled_rect(sld, 0, 0, 13.33, 7.5, DARK_BLUE)
    filled_rect(sld, 0, 5.6, 13.33, 1.9, MID_BLUE)
    tb(sld, title, 0.5, 1.5, 12.3, 2.5, size=36, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        tb(sld, subtitle, 0.5, 4.2, 12.3, 1.2, size=20,
           color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    tb(sld, "BS/C/7119 — Apply Fundamentals of Accounting\nFacilitator: S. N. Ngigi | Kenya School of TVET",
       0.5, 5.8, 12.3, 1.5, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    return sld

def section_divider(title):
    sld = prs.slides.add_slide(BLANK)
    filled_rect(sld, 0, 0, 13.33, 7.5, MID_BLUE)
    filled_rect(sld, 0, 3.2, 13.33, 0.08, ORANGE)
    tb(sld, title, 0.5, 2.2, 12.3, 2.0, size=32, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    return sld

def content_slide(heading, bullets, notes=""):
    sld = prs.slides.add_slide(BLANK)
    filled_rect(sld, 0, 0, 13.33, 1.1, DARK_BLUE)
    filled_rect(sld, 0, 1.1, 0.08, 6.4, ORANGE)
    tb(sld, heading, 0.3, 0.15, 12.7, 0.85, size=22, bold=True,
       color=WHITE, align=PP_ALIGN.LEFT)
    tb_lines(sld, bullets, 0.35, 1.25, 12.7, 5.9, size=17, color=DARK_GREY,
             line_space=4)
    return sld

def two_col_slide(heading, left_lines, right_lines, left_head="", right_head=""):
    sld = prs.slides.add_slide(BLANK)
    filled_rect(sld, 0, 0, 13.33, 1.1, DARK_BLUE)
    filled_rect(sld, 0, 1.1, 0.08, 6.4, ORANGE)
    tb(sld, heading, 0.3, 0.15, 12.7, 0.85, size=22, bold=True,
       color=WHITE, align=PP_ALIGN.LEFT)
    filled_rect(sld, 0.35, 1.25, 6.1, 5.9, RGBColor(0xF2, 0xF7, 0xFD))
    filled_rect(sld, 6.9, 1.25, 6.1, 5.9, RGBColor(0xF2, 0xF7, 0xFD))
    if left_head:
        tb(sld, left_head, 0.45, 1.3, 5.9, 0.45, size=15, bold=True,
           color=MID_BLUE)
    if right_head:
        tb(sld, right_head, 7.0, 1.3, 5.9, 0.45, size=15, bold=True,
           color=MID_BLUE)
    top = 1.3 + (0.45 if left_head else 0)
    tb_lines(sld, left_lines,  0.45, top, 5.9, 5.5, size=15, color=DARK_GREY)
    tb_lines(sld, right_lines, 7.0,  top, 5.9, 5.5, size=15, color=DARK_GREY)
    return sld

def table_slide(heading, col_headers, rows, col_widths=None):
    sld = prs.slides.add_slide(BLANK)
    filled_rect(sld, 0, 0, 13.33, 1.1, DARK_BLUE)
    filled_rect(sld, 0, 1.1, 0.08, 6.4, ORANGE)
    tb(sld, heading, 0.3, 0.15, 12.7, 0.85, size=22, bold=True,
       color=WHITE, align=PP_ALIGN.LEFT)
    cols = len(col_headers)
    if col_widths is None:
        col_widths = [12.6 / cols] * cols
    row_h = Inches(0.42)
    tbl_rows = 1 + len(rows)
    tbl = sld.shapes.add_table(tbl_rows, cols,
                                Inches(0.35), Inches(1.3),
                                Inches(12.6), row_h * tbl_rows).table
    # header row
    for ci, hdr in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = MID_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = WHITE
    # data rows
    for ri, row in enumerate(rows):
        bg = RGBColor(0xF2, 0xF7, 0xFD) if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_GREY
    # set col widths
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)
    return sld

def code_slide(heading, code_text):
    sld = prs.slides.add_slide(BLANK)
    filled_rect(sld, 0, 0, 13.33, 1.1, DARK_BLUE)
    filled_rect(sld, 0, 1.1, 0.08, 6.4, ORANGE)
    tb(sld, heading, 0.3, 0.15, 12.7, 0.85, size=22, bold=True,
       color=WHITE, align=PP_ALIGN.LEFT)
    filled_rect(sld, 0.35, 1.25, 12.6, 5.9, RGBColor(0x1E, 0x1E, 0x1E))
    tb(sld, code_text, 0.55, 1.35, 12.2, 5.7, size=13,
       color=RGBColor(0xD4, 0xD4, 0xD4), bold=False)
    return sld


# ════════════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# ════════════════════════════════════════════════════════════════════════════

title_slide("WEEK 5: ACCOUNT FOR FIXED ASSETS\n(NON-CURRENT ASSETS)", "Unit Code: BS/C/7119")

content_slide("Learning Outcomes", [
    "By the end of this presentation, trainees should be able to:",
    "",
    "1.  Define the term fixed asset and explain its recognition under relevant accounting standards",
    "2.  Account for acquisition of fixed assets",
    "3.  Compute depreciation using different methods",
    "4.  Account for depreciation charges",
    "5.  Explain and account for the provision for depreciation",
    "6.  Account for the disposal of fixed assets",
    "7.  Calculate profit and loss on disposal of fixed assets, and record the necessary journal entries",
    "8.  Apply accounting principles in handling fixed assets",
])

# ── LO1 ──────────────────────────────────────────────────────────────────────
content_slide("LO1 — Definition of Fixed Asset", [
    "A fixed asset (non-current asset) is an asset acquired for long-term use in the",
    "business (more than one accounting period) and NOT intended for resale.",
    "",
    "Under IAS 16 — Property, Plant and Equipment, a fixed asset is a tangible item that:",
    "  ✔  Is held for use in production, supply of goods/services, rental, or administration; AND",
    "  ✔  Is expected to be used during MORE THAN ONE accounting period (IASB, 2003, para. 6)",
    "",
    "Examples (Kenya context):",
    "  • Land & buildings  •  Motor vehicles  •  Plant & machinery",
    "  • Furniture & fittings  •  Computer equipment",
])

content_slide("LO1 — Recognition Criteria (IAS 16, para. 7)", [
    "A fixed asset is recognised in the books ONLY when BOTH conditions are met:",
    "",
    "  1.  It is PROBABLE that future economic benefits will flow to the entity",
    "",
    "  2.  The COST of the asset can be MEASURED RELIABLY",
    "",
    "If either condition is NOT met → expenditure is treated as an EXPENSE",
    "in the period it is incurred (not capitalised).",
    "",
    "⚠  Note: Land is NOT depreciated — it has an unlimited useful life.",
    "   (Exception: quarries, leasehold land)",
])

# ── LO2 ──────────────────────────────────────────────────────────────────────
content_slide("LO2 — Acquisition of Fixed Assets", [
    "The acquisition of a fixed asset is the process by which an entity obtains",
    "ownership or control of a non-current asset — whether through:",
    "  • Outright purchase   • Construction   • Exchange   • Finance lease",
    "",
    "The asset is initially recorded at its COST (Weygandt, Kimmel, & Kieso, 2018).",
    "",
    "Historical Cost (IAS 16, para. 16) includes:",
    "  ✔  Purchase price (less trade discounts, plus import duties)",
    "  ✔  Directly attributable costs: delivery, installation, testing",
    "  ✘  NOT: routine repairs, admin overhead, costs after asset is operational",
])

content_slide("LO2 — Example: Cost of a Machine [QUESTION]", [
    "Mwangi is setting up a factory in Athi River.",
    "He incurs the following costs in acquiring a new machine:",
    "",
    "  • Purchase price:          KES 800,000",
    "  • Delivery charges:        KES  25,000",
    "  • Installation costs:      KES  40,000",
    "  • Testing and trial runs:  KES  15,000",
    "",
    "Required: Calculate the TOTAL COST to be capitalised for this machine.",
])

table_slide("LO2 — Cost of a Machine — ANSWER (Mwangi, Athi River)",
    ["Item", "Amount (KES)"],
    [["Purchase price","800,000"],["Delivery charges","25,000"],
     ["Installation costs","40,000"],["Testing and trial runs","15,000"],
     ["TOTAL COST TO CAPITALISE","880,000"]],
    col_widths=[8.5, 4.1])

table_slide("LO2 — Recording Acquisitions: Journal Entries",
    ["Method of Purchase", "Debit", "Credit"],
    [["Cash purchase","Fixed Asset Account","Cash / Bank"],
     ["Credit purchase","Fixed Asset Account","Supplier (Creditor)"],
     ["Loan-financed","Fixed Asset Account","Loan Account"]],
    col_widths=[4.2, 4.2, 4.2])

# ── LO3 ──────────────────────────────────────────────────────────────────────
content_slide("LO3 — Depreciation: Definition & Key Terms", [
    "Depreciation is the SYSTEMATIC ALLOCATION of the depreciable amount of a",
    "non-current asset over its estimated useful life (IASB, 2003, para. 6).",
    "",
    "  • Depreciable Amount  =  Cost − Residual (Salvage) Value",
    "  • Useful Life          =  Period over which the asset is expected to be used",
    "  • Net Book Value (NBV) =  Cost − Accumulated Depreciation",
    "",
    "Reasons: Wear & tear  |  Passage of time  |  Obsolescence  |  Depletion  |  Legal limits",
])

content_slide("LO3 — Straight-Line Method: Example [QUESTION]", [
    "Wanjiku Transport, Nairobi, purchases a delivery van:",
    "  • Cost: KES 2,000,000  |  Residual value: KES 200,000  |  Life: 5 years",
    "  • Method: Straight-Line",
    "",
    "Required:",
    "  (a) Calculate the annual depreciation charge",
    "  (b) Prepare a 5-year depreciation schedule",
])

content_slide("LO3 — Straight-Line Method: ANSWER", [
    "Formula:  Annual Depreciation = (Cost − Residual Value) ÷ Useful Life",
    "",
    "  = (2,000,000 − 200,000) ÷ 5  =  KES 360,000 per year",
])

table_slide("LO3 — Straight-Line Schedule (Wanjiku Transport)",
    ["Year","Cost (KES)","Annual Dep'n (KES)","Acc. Dep'n (KES)","NBV (KES)"],
    [["1","2,000,000","360,000","360,000","1,640,000"],
     ["2","2,000,000","360,000","720,000","1,280,000"],
     ["3","2,000,000","360,000","1,080,000","920,000"],
     ["4","2,000,000","360,000","1,440,000","560,000"],
     ["5","2,000,000","360,000","1,800,000","200,000 ✓"]],
    col_widths=[1.5, 2.5, 2.8, 2.8, 3.0])

content_slide("LO3 — Reducing Balance Method: Example [QUESTION]", [
    "Ochieng Engineering, Kisumu, purchases machinery:",
    "  • Cost: KES 1,000,000  |  Rate: 25% p.a. reducing balance",
    "",
    "Required:",
    "  (a) Calculate depreciation for each of the first 5 years",
    "  (b) Prepare a depreciation schedule showing NBV at end of each year",
])

content_slide("LO3 — Reducing Balance Method: ANSWER", [
    "Formula:  Annual Depreciation = Rate (%) × NBV at start of year",
    "",
    "Gives HIGHER depreciation in early years, LOWER in later years.",
    "The asset is NEVER fully depreciated to zero under this method.",
])

table_slide("LO3 — Reducing Balance Schedule (Ochieng Engineering)",
    ["Year","NBV at Start (KES)","Dep'n 25% (KES)","Acc. Dep'n (KES)","NBV at End (KES)"],
    [["1","1,000,000","250,000","250,000","750,000"],
     ["2","750,000","187,500","437,500","562,500"],
     ["3","562,500","140,625","578,125","421,875"],
     ["4","421,875","105,469","683,594","316,406"],
     ["5","316,406","79,102","762,696","237,304"]],
    col_widths=[1.5, 2.8, 2.5, 2.8, 3.0])

table_slide("LO3 — Comparison of Depreciation Methods",
    ["Feature","Straight-Line","Reducing Balance"],
    [["Annual charge","Equal each year","Decreasing each year"],
     ["Early years","Lower charge","Higher charge"],
     ["Later years","Same charge","Lower charge"],
     ["Best suited for","Furniture, buildings (even wear)","Vehicles, computers (rapid early loss)"],
     ["IAS 16 compliant","Yes","Yes"]],
    col_widths=[3.5, 4.5, 4.6])

# ── LO4 ──────────────────────────────────────────────────────────────────────
content_slide("LO4 — Accounting for Depreciation Charges", [
    "Each year, depreciation is recorded with the following double entry:",
    "",
    "  Dr  Depreciation Expense (Income Statement)   XXX",
    "  Cr  Provision for Depreciation Account            XXX",
    "",
    "Example — Wanjiku Transport (KES 360,000/year, straight-line):",
    "  Dr  Depreciation Expense          360,000",
    "  Cr  Provision for Depreciation              360,000",
    "",
    "At year end, Depreciation Expense is transferred to P&L:",
    "  Dr  Profit & Loss    360,000",
    "  Cr  Depreciation Expense          360,000",
])

# ── LO5 ──────────────────────────────────────────────────────────────────────
content_slide("LO5 — Provision for Depreciation: Concept", [
    "The Provision for Depreciation Account is a CONTRA-ASSET account that",
    "accumulates total depreciation charged on a fixed asset since acquisition.",
    "",
    "Why keep it separate from the asset account?",
    "  • The asset account stays at HISTORICAL COST throughout",
    "  • Accumulated depreciation is deducted on the Balance Sheet to show NBV",
    "",
    "Balance Sheet presentation:",
    "  Motor Vehicle (at cost)              KES 2,000,000",
    "  Less: Provision for Depreciation        (720,000)",
    "  Net Book Value                       KES 1,280,000",
])

content_slide("LO5 — Provision for Depreciation Ledger [QUESTION]", [
    "Using Wanjiku Transport's delivery van:",
    "  • Cost: KES 2,000,000  |  Annual depreciation: KES 360,000 (straight-line)",
    "",
    "Required — prepare ledger accounts for Years 1, 2, and 3:",
    "  (a) Motor Vehicle Account (at cost)",
    "  (b) Provision for Depreciation — Motor Vehicle Account",
    "  (c) Depreciation Expense Account",
])

code_slide("LO5 — Provision for Depreciation Ledger — ANSWER",
"""  Motor Vehicle Account (at cost) — remains at KES 2,000,000 throughout

  Provision for Depreciation — Motor Vehicle
  ──────────────────────────────────────────────────────────────
  Dr                                  Cr
  ──────────────────────────────────────────────────────────────
  YEAR 1:  Bal c/d    360,000  |  Dep'n Expense      360,000
  YEAR 2:  Bal c/d    720,000  |  Bal b/d    360,000
                               |  Dep'n Exp  360,000  = 720,000
  YEAR 3:  Bal c/d  1,080,000  |  Bal b/d    720,000
                               |  Dep'n Exp  360,000  = 1,080,000

  Depreciation Expense Account (each year):
  Dr  Prov for Dep'n  360,000  |  Cr  P&L  360,000
""")

# ── LO6 ──────────────────────────────────────────────────────────────────────
content_slide("LO6 — Disposal of Fixed Assets: When & Steps", [
    "Assets are disposed of when: sold  |  traded in  |  scrapped  |  stolen/destroyed",
    "",
    "A DISPOSAL ACCOUNT is opened. Steps:",
    "  1. Remove asset at cost       →  Dr Disposal A/c,  Cr Asset A/c",
    "  2. Remove accumulated dep'n   →  Dr Prov for Dep'n,  Cr Disposal A/c",
    "  3. Record proceeds            →  Dr Cash/Bank,  Cr Disposal A/c",
    "  4. Transfer balance:",
    "       Profit (credit balance)  →  Dr Disposal A/c,  Cr P&L",
    "       Loss   (debit balance)   →  Dr P&L,  Cr Disposal A/c",
])

content_slide("LO6 — Disposal at a PROFIT [QUESTION] (Kamau Traders, Nakuru)", [
    "Kamau Traders purchased office furniture on 1 January 2023 for KES 200,000.",
    "Depreciation: 10% per annum, straight-line.",
    "Sold on 1 January 2026 for KES 160,000 cash.",
    "",
    "Required:",
    "  (a) Calculate accumulated depreciation and NBV at disposal",
    "  (b) Prepare the Disposal Account",
])

content_slide("LO6 — Disposal at a PROFIT — ANSWER (Kamau Traders)", [
    "Accumulated Dep'n = 3 years × 10% × 200,000 = KES 60,000",
    "NBV at disposal = 200,000 − 60,000 = KES 140,000",
    "",
    "Disposal Account:",
    "  Dr  Furniture (cost)    200,000  |  Cr  Prov for Dep'n   60,000",
    "                                   |  Cr  Cash            160,000",
    "                                   |  Cr  P&L (Profit)     20,000",
    "  ─────────────────────────────────────────────────────────────────",
    "                          200,000  |                      200,000",
    "",
    "PROFIT ON DISPOSAL = KES 20,000  →  Income in P&L",
])

content_slide("LO6 — Disposal at a LOSS [QUESTION] (Atieno Supplies, Kisumu)", [
    "Atieno Supplies purchased a computer on 1 January 2024 for KES 120,000.",
    "Depreciation: 25% per annum, reducing balance.",
    "Sold on 31 December 2025 for KES 50,000.",
    "",
    "Required:",
    "  (a) Calculate Year 1 and Year 2 depreciation, accumulated dep'n and NBV",
    "  (b) Prepare the Disposal Account",
])

content_slide("LO6 — Disposal at a LOSS — ANSWER (Atieno Supplies)", [
    "Year 1: 25% × 120,000 = 30,000  →  NBV = 90,000",
    "Year 2: 25% × 90,000  = 22,500  →  NBV = 67,500",
    "Accumulated Dep'n = KES 52,500",
    "",
    "Disposal Account:",
    "  Dr  Computer (cost)  120,000  |  Cr  Prov for Dep'n   52,500",
    "                                |  Cr  Cash             50,000",
    "                                |  Cr  P&L (Loss)       17,500",
    "  ──────────────────────────────────────────────────────────────",
    "                       120,000  |                      120,000",
    "",
    "LOSS ON DISPOSAL = KES 17,500  →  Expense in P&L",
])

# ── LO7 ──────────────────────────────────────────────────────────────────────
table_slide("LO7 — Journal Entries for Disposal",
    ["Step","Narrative","Debit","Credit"],
    [["1","Remove asset at cost","Disposal Account","Asset Account (cost)"],
     ["2","Remove accumulated dep'n","Provision for Depreciation","Disposal Account"],
     ["3","Record proceeds","Cash / Bank","Disposal Account"],
     ["4a","Profit (credit balance)","Disposal Account","P&L — Profit on Disposal"],
     ["4b","Loss (debit balance)","P&L — Loss on Disposal","Disposal Account"]],
    col_widths=[1.0, 3.8, 3.8, 4.0])

table_slide("LO7 — Journal Entries: Kamau Traders (Profit Example)",
    ["Step","Account","Dr (KES)","Cr (KES)"],
    [["1","Disposal Account","200,000",""],
     ["1","Furniture Account (at cost)","","200,000"],
     ["2","Provision for Depreciation","60,000",""],
     ["2","Disposal Account","","60,000"],
     ["3","Cash / Bank","160,000",""],
     ["3","Disposal Account","","160,000"],
     ["4","Disposal Account","20,000",""],
     ["4","P&L — Profit on Disposal","","20,000"]],
    col_widths=[0.8, 6.5, 2.6, 2.7])

table_slide("LO7 — Journal Entries: Atieno Supplies (Loss Example)",
    ["Step","Account","Dr (KES)","Cr (KES)"],
    [["1","Disposal Account","120,000",""],
     ["1","Computer Account (at cost)","","120,000"],
     ["2","Provision for Depreciation","52,500",""],
     ["2","Disposal Account","","52,500"],
     ["3","Cash / Bank","50,000",""],
     ["3","Disposal Account","","50,000"],
     ["4","P&L — Loss on Disposal","17,500",""],
     ["4","Disposal Account","","17,500"]],
    col_widths=[0.8, 6.5, 2.6, 2.7])

# ── LO8 ──────────────────────────────────────────────────────────────────────
table_slide("LO8 — Accounting Principles in Handling Fixed Assets",
    ["Principle / Standard","Application to Fixed Assets"],
    [["Historical Cost (IAS 16)","Initially recorded at cost — amount paid to acquire & bring into use"],
     ["Matching (Accruals)","Depreciation charged in the period the asset generates revenue"],
     ["Consistency (IAS 16, para. 61)","Same depreciation method applied each period; changes require disclosure"],
     ["Prudence","Assets not overstated; impairment loss recognised if NBV > recoverable amount (IAS 36)"],
     ["Going Concern","Assets depreciated over useful life assuming business will continue"],
     ["Materiality","Low-value items may be expensed immediately rather than capitalised"],
     ["Substance over Form (IFRS 16)","Finance lease assets capitalised even without legal ownership"]],
    col_widths=[4.5, 8.1])

content_slide("LO8 — IAS 16 Summary", [
    "Initial recognition:    at COST",
    "",
    "Subsequent measurement — choose one model:",
    "  • Cost Model:        Cost less accumulated depreciation and impairment losses",
    "  • Revaluation Model: Fair value less subsequent depreciation and impairment",
    "",
    "Depreciation method must reflect the PATTERN OF CONSUMPTION of economic benefits.",
    "Residual value & useful life reviewed at least at EACH FINANCIAL YEAR-END.",
])

# ── PRACTICE QUESTIONS ───────────────────────────────────────────────────────
content_slide("Practice Q1 — Depreciation [QUESTION]", [
    "Nyambura Agribusiness, Nanyuki, purchased a tractor on 1 July 2024",
    "for KES 3,600,000. Residual value: KES 600,000. Useful life: 5 years. Year end: 30 June.",
    "",
    "Required:",
    "  (a) Depreciation schedule for 5 years — STRAIGHT-LINE method",
    "  (b) Depreciation schedule for 5 years — REDUCING BALANCE at 30% p.a.",
])

table_slide("Q1(a) ANSWER — Straight-Line  |  Annual Dep'n = (3,600,000−600,000)÷5 = KES 600,000",
    ["Year Ending","Cost (KES)","Annual Dep'n (KES)","Acc. Dep'n (KES)","NBV (KES)"],
    [["30 Jun 2025","3,600,000","600,000","600,000","3,000,000"],
     ["30 Jun 2026","3,600,000","600,000","1,200,000","2,400,000"],
     ["30 Jun 2027","3,600,000","600,000","1,800,000","1,800,000"],
     ["30 Jun 2028","3,600,000","600,000","2,400,000","1,200,000"],
     ["30 Jun 2029","3,600,000","600,000","3,000,000","600,000 ✓"]],
    col_widths=[2.2, 2.3, 2.8, 2.8, 2.5])

table_slide("Q1(b) ANSWER — Reducing Balance at 30% p.a.",
    ["Year Ending","NBV at Start (KES)","Dep'n 30% (KES)","Acc. Dep'n (KES)","NBV at End (KES)"],
    [["30 Jun 2025","3,600,000","1,080,000","1,080,000","2,520,000"],
     ["30 Jun 2026","2,520,000","756,000","1,836,000","1,764,000"],
     ["30 Jun 2027","1,764,000","529,200","2,365,200","1,234,800"],
     ["30 Jun 2028","1,234,800","370,440","2,735,640","864,360"],
     ["30 Jun 2029","864,360","259,308","2,994,948","605,052"]],
    col_widths=[2.2, 2.6, 2.5, 2.8, 2.5])

content_slide("Practice Q2 — Disposal [QUESTION]", [
    "Mutua Logistics, Machakos, purchased a lorry on 1 January 2023 for KES 4,000,000.",
    "Depreciation: 20% straight-line, no residual value.",
    "Lorry sold on 30 June 2025 for KES 2,200,000.",
    "",
    "Required:",
    "  (a) Calculate accumulated depreciation at the date of disposal",
    "  (b) Prepare the Disposal Account and all journal entries",
    "  (c) State whether there was a profit or loss on disposal",
])

content_slide("Q2(a) ANSWER — Accumulated Depreciation", [
    "Annual depreciation = 20% × KES 4,000,000 = KES 800,000",
    "",
    "  2023 (full year):              KES 800,000",
    "  2024 (full year):              KES 800,000",
    "  2025 (half year — Jan–Jun):    KES 400,000",
    "  ─────────────────────────────────────────",
    "  Total accumulated depreciation: KES 2,000,000",
    "",
    "  NBV at disposal = 4,000,000 − 2,000,000 = KES 2,000,000",
])

table_slide("Q2(b) ANSWER — Journal Entries",
    ["Step","Account","Dr (KES)","Cr (KES)"],
    [["1","Disposal Account","4,000,000",""],
     ["1","Lorry Account (at cost)","","4,000,000"],
     ["2","Provision for Depreciation — Lorry","2,000,000",""],
     ["2","Disposal Account","","2,000,000"],
     ["3","Cash / Bank","2,200,000",""],
     ["3","Disposal Account","","2,200,000"],
     ["4","Disposal Account","200,000",""],
     ["4","P&L — Profit on Disposal","","200,000"]],
    col_widths=[0.8, 6.5, 2.6, 2.7])

code_slide("Q2(b) Disposal Account  |  Q2(c): PROFIT of KES 200,000",
"""  Lorry Disposal Account
  ──────────────────────────────────────────────────────────────────
  Dr                                       Cr
  ──────────────────────────────────────────────────────────────────
  Lorry (cost)       4,000,000  |  Prov for Dep'n      2,000,000
  P&L (Profit)         200,000  |  Cash                2,200,000
  ──────────────────────────────────────────────────────────────────
                     4,200,000  |                      4,200,000

  ✔  PROFIT ON DISPOSAL = KES 200,000  →  Income in Statement of P&L
""")

content_slide("Practice Q3 — Discussion [QUESTION]", [
    "(a) Explain why land is generally NOT depreciated but buildings ARE.",
    "",
    "(b) A Nairobi business bought computers in 2020. Explain why the",
    "    REDUCING BALANCE method may be more appropriate than straight-line.",
    "",
    "(c) State TWO accounting principles that support depreciating fixed assets.",
])

content_slide("Q3 — ANSWERS", [
    "(a) Land has an UNLIMITED useful life — does not wear out or become obsolete.",
    "    Buildings have a FINITE useful life (wear & tear, weather) → depreciated (IAS 16).",
    "",
    "(b) Computers lose value RAPIDLY early due to technological obsolescence.",
    "    Reducing balance charges HIGHER depreciation early, matching actual value loss.",
    "    Also matches higher maintenance costs in later years with lower dep'n charges.",
    "",
    "(c) 1. Matching Principle — matches asset cost to revenue earned each period.",
    "    2. Prudence — prevents overstatement of asset values and profit.",
])

# ── REFERENCES & END ─────────────────────────────────────────────────────────
content_slide("References", [
    "International Accounting Standards Board (IASB). (2003).",
    "  IAS 16: Property, plant and equipment (revised 2003). IFRS Foundation.",
    "  https://www.ifrs.org/issued-standards/list-of-standards/ias-16-property-plant-and-equipment/",
    "",
    "Weygandt, J. J., Kimmel, P. D., & Kieso, D. E. (2018).",
    "  Financial accounting: IFRS edition (3rd ed.). John Wiley & Sons.",
])

title_slide("End of Week 5", "Fixed Assets & Depreciation")

out = "/home/stdk/NOTES_KSTVET/Fundamentals of Accountig _BS_C_7119/Week_5_Fixed_Assets_and_Depreciation.pptx"
prs.save(out)
print(f"Saved: {out}")

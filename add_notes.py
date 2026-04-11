from pptx import Presentation

PATH = "/home/stdk/NOTES_KSTVET/Fundamentals of Accountig _BS_C_7119/Week_5_Fixed_Assets_and_Depreciation.pptx"
prs = Presentation(PATH)
s = prs.slides  # shorthand

def note(idx, text):
    s[idx].notes_slide.notes_text_frame.text = text.strip()

note(0, """
Welcome trainees to Week 5 of BS/C/7119 — Apply Fundamentals of Accounting.
This week we focus on fixed assets (non-current assets): how to define, acquire, depreciate, and dispose of them.
By the end of this session you should be able to handle all aspects of fixed asset accounting as required by IAS 16.
Remind trainees to have their exercise books ready for worked examples.
""")

note(1, """
Walk through each learning outcome clearly. Emphasise that LOs 1–8 build on each other:
- You must understand what a fixed asset IS (LO1) before you can record its acquisition (LO2).
- You must know how to compute depreciation (LO3) before you can account for it (LO4) or maintain the provision account (LO5).
- Disposal (LO6 & LO7) requires knowledge of both cost and accumulated depreciation.
- LO8 ties everything together with the underlying accounting principles.
Ask trainees: "Can anyone give an example of a fixed asset used in a Kenyan business?"
""")

note(2, """
Key teaching points:
1. Stress the phrase "more than one accounting period" — this is what distinguishes a fixed asset from a current asset or an expense.
2. IAS 16 is the governing standard. Kenya adopted IFRS for SMEs and full IFRS, so IAS 16 applies.
3. The two-part IAS 16 definition (held for use AND more than one period) is examinable — trainees should memorise both parts.
4. Go through the Kenya examples: ask trainees to suggest additional examples from their own experience (e.g., a boda boda, a posho mill, a school building).
5. Clarify: a matatu bought for resale by a dealer is STOCK (current asset), not a fixed asset. Intent of use matters.
""")

note(3, """
Recognition criteria are critical — many exam questions test whether an item should be capitalised or expensed.
Walk through both conditions:
  Condition 1 — Probable future economic benefits: the asset must be expected to generate income or reduce costs.
  Condition 2 — Reliable measurement of cost: there must be a verifiable invoice or valuation.

Practical example: A business pays a deposit for a machine not yet delivered. Should it be recognised?
  → Condition 2 may not yet be fully met if the final cost is uncertain. Discuss with trainees.

Land note: Reinforce that land is NOT depreciated. This is a common exam error. Exception: leasehold land (finite legal right) and quarries (depleted over time).
""")

note(4, """
Acquisition definition (Weygandt et al., 2018): stress that acquisition is not just buying — it includes construction, exchange, and finance leases.
Historical cost principle (IAS 16, para. 16): the asset is recorded at what it COST to get it ready for use.

Spend time on the "directly attributable costs" concept:
  ✔ Delivery, installation, testing — these are necessary to bring the asset into working condition.
  ✘ Repairs after the asset is in use, admin costs, training staff to use the asset — these are EXPENSES.

Ask trainees: "If a business pays KES 10,000 to train staff to use a new machine, should this be capitalised?"
  → No. Training costs are not directly attributable to bringing the asset to its working condition.
""")

note(5, """
Before revealing the answer, give trainees 2–3 minutes to attempt this individually.
Remind them: only costs that are DIRECTLY ATTRIBUTABLE to bringing the asset to its present location and condition are capitalised.
Common mistake: students sometimes include insurance or finance charges — clarify these are NOT capitalised under IAS 16.
""")

note(6, """
Walk through each line:
- Purchase price: the invoice amount — always capitalised.
- Delivery: directly attributable — capitalised.
- Installation: directly attributable — capitalised.
- Testing: directly attributable (asset not yet operational) — capitalised.
Total = KES 880,000. This is the amount entered in the asset account.

Ask: "What if Mwangi also paid KES 5,000 for a maintenance contract for the first year?"
  → NOT capitalised. Maintenance is an ongoing expense, not part of bringing the asset into use.
""")

note(7, """
Three methods of purchase — all result in the same debit to the Fixed Asset Account at cost.
The credit differs depending on how payment is made.

Walk through each:
1. Cash: straightforward — asset bought and paid immediately.
2. Credit: asset received but not yet paid — creates a creditor (liability).
3. Loan: asset paid for by borrowing — creates a loan liability.

Emphasise: regardless of payment method, the asset is ALWAYS recorded at its full cost on the date of acquisition.
Ask trainees to write the journal entry for Mwangi's machine (KES 880,000) if purchased on credit from Nairobi Machinery Ltd.
""")

note(8, """
Depreciation is one of the most important concepts in this unit.
Key points to emphasise:
1. Depreciation is an ALLOCATION of cost — not a valuation exercise. We are not trying to show the market value of the asset.
2. The depreciable amount = Cost minus Residual Value. If residual value is zero, the full cost is depreciated.
3. NBV (Net Book Value) = what the asset is shown at on the Balance Sheet at any point in time.

Reasons for depreciation — use relatable Kenyan examples:
  - Wear & tear: a matatu engine wears out from daily use on Nairobi roads.
  - Obsolescence: a 2010 computer is outdated by 2024 models.
  - Depletion: a quarry in Machakos is gradually exhausted.

Ask: "Why don't we just expense the full cost of a machine in the year we buy it?"
  → Because it generates income over many years — the Matching Principle requires us to spread the cost.
""")

note(9, """
Give trainees 3–4 minutes to attempt parts (a) and (b) before revealing the answer.
Remind them of the formula: Annual Depreciation = (Cost − Residual Value) ÷ Useful Life.
Encourage them to set up the schedule in a table format in their exercise books.
""")

note(10, """
Walk through the formula step by step:
  (2,000,000 − 200,000) ÷ 5 = 1,800,000 ÷ 5 = KES 360,000 per year.

Key features of straight-line:
- The SAME amount is charged every year.
- At the end of Year 5, NBV = Residual Value (KES 200,000). This is a useful check.
- The asset account stays at KES 2,000,000 throughout — only the provision account changes.

Journal entry: Dr Depreciation Expense / Cr Provision for Depreciation.
Ask: "Where does the Depreciation Expense appear in the financial statements?"
  → Income Statement (reduces profit). The Provision appears on the Balance Sheet (deducted from asset cost).
""")

note(11, """
Go through the schedule row by row. Highlight:
- Cost column never changes — always KES 2,000,000 (historical cost principle).
- Annual depreciation is constant — KES 360,000 every year.
- Accumulated depreciation grows by KES 360,000 each year.
- NBV decreases by KES 360,000 each year.
- At Year 5: NBV = KES 200,000 = Residual Value ✓ (good check for trainees).

Ask trainees to verify: 2,000,000 − 1,800,000 = 200,000. ✓
""")

note(12, """
Give trainees 3–4 minutes to attempt parts (a) and (b).
Remind them: under reducing balance, the rate is applied to the NBV at the START of each year, not the original cost.
This is the key difference from straight-line.
""")

note(13, """
Key features of reducing balance:
- Higher depreciation in early years, lower in later years.
- Reflects the reality that many assets (vehicles, computers) lose value faster when new.
- The asset is NEVER fully depreciated to zero — it asymptotically approaches zero.
- No residual value is needed in the formula (the method naturally leaves a residual).

Ask: "Which method gives a higher depreciation charge in Year 1 for the same asset?"
  → Reducing balance (for the same asset, RBM Year 1 > SLM Year 1 in most cases).
""")

note(14, """
Go through the schedule. Highlight:
- NBV at Start of Year 2 = NBV at End of Year 1 (KES 750,000).
- Each year's depreciation is smaller than the previous year.
- After 5 years, NBV = KES 237,304 — the asset still has a book value.

Compare with the SLM schedule for Wanjiku's van:
  SLM Year 1 dep'n: KES 360,000
  RBM Year 1 dep'n (if 25% on KES 2m): KES 500,000 — higher under RBM.
""")

note(15, """
Use this comparison table to consolidate understanding of both methods.
Key exam point: IAS 16 allows BOTH methods — the choice should reflect the pattern of economic benefit consumption.
- Straight-line: best for assets that provide equal benefit each year (furniture, buildings).
- Reducing balance: best for assets that provide more benefit early on (vehicles, computers, technology).

Consistency principle: once a method is chosen, it must be applied consistently. A change requires disclosure in the financial statements (IAS 16, para. 61).

Ask: "A business changes from straight-line to reducing balance for its vehicles. What must it disclose?"
  → The reason for the change, and the effect on the current period's depreciation charge.
""")

note(16, """
This slide covers the double entry for recording depreciation each year.
Walk through both entries:
1. Dr Depreciation Expense / Cr Provision for Depreciation — records the annual charge.
2. Dr P&L / Cr Depreciation Expense — closes the expense account to the Income Statement at year end.

Emphasise: the ASSET ACCOUNT IS NEVER TOUCHED. It stays at historical cost.
The provision account accumulates on the credit side — it grows each year.

Ask: "After 3 years, what is the balance on the Provision for Depreciation account for Wanjiku's van?"
  → 3 × 360,000 = KES 1,080,000 (credit balance).
""")

note(17, """
The provision for depreciation account is a CONTRA-ASSET — it has a credit balance and is deducted from the asset on the Balance Sheet.

Balance Sheet presentation is important for exams:
  Motor Vehicle (at cost)              2,000,000
  Less: Provision for Depreciation      (720,000)   ← after Year 2
  Net Book Value                       1,280,000

Ask trainees: "Why do we keep the asset at cost and use a separate provision account, rather than just reducing the asset account directly?"
  → So users of financial statements can see both the original cost AND the accumulated depreciation. This gives more information about the age and condition of assets.
""")

note(18, """
Give trainees 4–5 minutes to attempt all three ledger accounts.
Remind them of the format: Dr side on the left, Cr side on the right, with dates and narrations.
Common mistakes:
- Putting depreciation on the wrong side of the provision account (it should be on the CREDIT side).
- Forgetting to bring down the balance (Bal b/d) at the start of each new year.
""")

note(19, """
Walk through each account carefully:
1. Motor Vehicle Account: stays at KES 2,000,000 on the debit side throughout. Never changes.
2. Provision for Depreciation: credit side grows each year. Balance b/d + new charge = new balance c/d.
   Year 1: 360,000. Year 2: 720,000. Year 3: 1,080,000.
3. Depreciation Expense: opened each year, debited with the charge, then closed to P&L (credited).
   It has a nil balance at year end after the transfer.

Ask: "What is the NBV of the van at the end of Year 3?"
  → 2,000,000 − 1,080,000 = KES 920,000.
""")

note(20, """
Disposal occurs when an asset is no longer used by the business.
The key tool is the DISPOSAL ACCOUNT — a temporary account opened to calculate profit or loss.

Walk through the 4 steps carefully. Stress:
- Step 1: We REMOVE the asset from the books at its original cost (debit disposal, credit asset).
- Step 2: We REMOVE the accumulated depreciation (debit provision, credit disposal).
  After steps 1 & 2, the disposal account shows the NBV on the debit side.
- Step 3: We record what we received (debit cash, credit disposal).
- Step 4: The balance is profit (credit balance → income) or loss (debit balance → expense).

Ask: "After steps 1 and 2, what does the balance on the disposal account represent?"
  → The Net Book Value of the asset at the date of disposal.
""")

note(21, """
Give trainees 4–5 minutes to attempt all five parts.
Remind them to work through the steps in order:
1. Calculate accumulated depreciation first.
2. Then NBV.
3. Compare NBV with proceeds to find profit or loss.
4. Then prepare the disposal account.
5. Then write the journal entries.
""")

note(22, """
Walk through the solution step by step:
- Accumulated dep'n: 3 years × 10% × 200,000 = KES 60,000. (Straight-line: same amount each year.)
- NBV: 200,000 − 60,000 = KES 140,000.
- Proceeds: KES 160,000 > NBV KES 140,000 → PROFIT of KES 20,000.

Disposal account: debit side = cost (200,000). Credit side = provision (60,000) + cash (160,000) + profit (20,000) = 240,000.
Wait — let's recheck: 60,000 + 160,000 = 220,000. Debit = 200,000. Credit exceeds debit by 20,000 → profit.
The profit entry: Dr Disposal A/c 20,000 / Cr P&L 20,000. This balances the account.

Ask: "Where does the profit on disposal appear in the financial statements?"
  → Income Statement — as other income (not part of trading profit).
""")

note(23, """
Give trainees 4–5 minutes to attempt all five parts.
Remind them that under reducing balance, they must calculate Year 1 depreciation first, find the new NBV, then calculate Year 2 depreciation on that NBV.
""")

note(24, """
Walk through:
- Year 1: 25% × 120,000 = 30,000. NBV = 90,000.
- Year 2: 25% × 90,000 = 22,500. NBV = 67,500.
- Accumulated dep'n = 52,500.
- Proceeds = 50,000 < NBV 67,500 → LOSS of 17,500.

Disposal account: debit = cost (120,000). Credit = provision (52,500) + cash (50,000) = 102,500.
Debit exceeds credit by 17,500 → LOSS. Entry: Dr P&L 17,500 / Cr Disposal A/c 17,500.

Ask: "Where does the loss on disposal appear in the financial statements?"
  → Income Statement — as an expense (increases the loss or reduces profit).
""")

note(25, """
This slide summarises the journal entry template for ALL disposals.
Trainees should memorise the 4 steps — they are always the same regardless of whether there is a profit or loss.
The only difference is Step 4: profit goes to P&L as income, loss goes to P&L as expense.

Tip for exams: always open the disposal account first, fill in steps 1–3, then find the balancing figure (profit or loss) before writing the journal entries.
""")

note(26, """
Walk through each journal entry for Kamau Traders line by line.
Stress that each step has TWO entries (debit and credit) — this is double-entry bookkeeping.
After all 4 steps, the disposal account balances to zero.

Ask trainees to total both sides of the disposal account:
  Debit: 200,000 + 20,000 = 220,000.
  Credit: 60,000 + 160,000 = 220,000. ✓ Balanced.
""")

note(27, """
Walk through each journal entry for Atieno Supplies.
After all 4 steps, the disposal account should balance:
  Debit: 120,000.
  Credit: 52,500 + 50,000 + 17,500 = 120,000. ✓ Balanced.

Reinforce: the loss (KES 17,500) is debited to P&L — it increases expenses and reduces profit.
Ask: "How would this loss affect the business's tax liability?"
  → A loss on disposal is a deductible expense — it reduces taxable profit (subject to tax rules).
""")

note(28, """
This slide ties the entire week's content to the underlying accounting principles.
Go through each principle and link it to a specific example from the week:
- Historical Cost: Mwangi's machine recorded at KES 880,000 (not market value).
- Matching: Wanjiku's van depreciated KES 360,000/year — matched to the years it generates income.
- Consistency: Ochieng must use 25% RBM every year unless there is a justified reason to change.
- Prudence: If Wanjiku's van is worth less than its NBV, an impairment loss must be recognised (IAS 36).
- Going Concern: We depreciate over 5 years assuming the business will still be operating in Year 5.
- Materiality: A KES 500 stapler need not be capitalised — expense it immediately.
- Substance over Form: A finance lease asset is capitalised even if the business doesn't legally own it yet.
""")

note(29, """
IAS 16 summary — key points for exams:
1. Initial recognition: ALWAYS at cost. Never at estimated value or market price.
2. Two subsequent measurement models: cost model (most common) and revaluation model.
3. Depreciation method must reflect the pattern of consumption — this is why the choice of method matters.
4. Residual value and useful life must be reviewed at EVERY year end — if circumstances change, the depreciation charge changes (prospectively, not retrospectively).

Ask: "A business originally estimated a machine's useful life as 10 years. After 3 years, it revises this to 7 years total. How does this affect depreciation?"
  → Remaining depreciable amount is spread over the remaining 4 years (7 − 3). This is a change in accounting estimate, not a change in accounting policy.
""")

note(30, """
Practice Question 1 — give trainees 8–10 minutes to attempt both parts.
Remind them:
- SLM: Annual dep'n = (3,600,000 − 600,000) ÷ 5 = KES 600,000.
- RBM: Year 1 = 30% × 3,600,000. Year 2 = 30% × NBV at end of Year 1. And so on.
Encourage trainees to set up the table before filling in numbers.
""")

note(31, """
SLM answer: Annual depreciation = KES 600,000 every year.
At Year 5: NBV = KES 600,000 = Residual Value ✓.
Ask trainees to check their own work against this table.
Note: the cost column never changes (KES 3,600,000 throughout).
""")

note(32, """
RBM answer: depreciation decreases each year as the NBV falls.
Note that after 5 years, NBV = KES 605,052 — slightly above the residual value of KES 600,000.
Under RBM, the asset is never fully depreciated to the residual value unless the rate is chosen precisely.
Ask: "Which method gives a higher total depreciation charge over 5 years?"
  → Both methods depreciate the same total amount (cost minus residual value = KES 3,000,000) over the asset's life. The difference is only in the TIMING of the charges.
""")

note(33, """
Practice Question 2 — give trainees 10 minutes to attempt all three parts.
Key challenge: time apportionment. The lorry is sold on 30 June 2025 — only 6 months into 2025.
Remind trainees: depreciation for 2025 = 20% × 4,000,000 × 6/12 = KES 400,000.
""")

note(34, """
Walk through the time apportionment carefully:
- 2023: full year = KES 800,000.
- 2024: full year = KES 800,000.
- 2025: half year (Jan–Jun) = KES 400,000.
- Total = KES 2,000,000.
NBV at disposal = 4,000,000 − 2,000,000 = KES 2,000,000.
Ask: "Why do we only charge half a year's depreciation in 2025?"
  → The asset was only used for 6 months of 2025 before disposal. Matching principle.
""")

note(35, """
Walk through each journal entry. After all 4 steps:
  Debit side of disposal account: 4,000,000 + 200,000 = 4,200,000.
  Credit side: 2,000,000 + 2,200,000 = 4,200,000. ✓ Balanced.
The profit of KES 200,000 arises because the lorry was sold for MORE than its NBV.
""")

note(36, """
The disposal account confirms the profit of KES 200,000.
Credit side (4,200,000) exceeds debit side (4,000,000) by KES 200,000 → PROFIT.
This profit is recorded as income in the Statement of Profit or Loss for the period ending 30 June 2025.
Ask: "Would this profit be included in gross profit or below the line?"
  → Below the line — it is not part of normal trading activity. It appears as 'other income'.
""")

note(37, """
Discussion questions — give trainees 5 minutes to think and write answers before discussing as a class.
These questions test conceptual understanding, not just computation.
Encourage trainees to use accounting terminology in their answers.
""")

note(38, """
Model answers — go through each part:
(a) Land: unlimited life → no depreciation. Buildings: finite life → depreciated. IAS 16 requires depreciation only where useful life is finite.
(b) Computers: rapid early obsolescence → reducing balance better reflects the pattern of economic benefit consumption. Also matches higher maintenance costs in later years with lower depreciation charges.
(c) Matching Principle and Prudence — both are examinable. Trainees should be able to explain HOW each principle applies, not just name it.
""")

note(39, """
References slide — remind trainees of the importance of citing sources in their written work.
IAS 16 is the primary standard for fixed assets. Available free on the IFRS Foundation website.
Weygandt et al. is a widely used financial accounting textbook — available in the library.
For TVET assessments, trainees are expected to reference IAS 16 when explaining recognition and measurement of fixed assets.
""")

note(40, """
End of Week 5.
Recap the 8 learning outcomes — ask trainees to self-assess: which LOs are they confident about? Which need more practice?
Assign the Standard Assessment Questions document as homework — all 12 questions to be attempted before the next session.
Next week: [insert Week 6 topic].
Thank trainees for their participation.
""")

prs.save(PATH)
print(f"Speaker notes added and saved: {PATH}")
